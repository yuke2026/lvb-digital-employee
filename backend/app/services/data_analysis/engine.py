"""
数据经营分析引擎
输入：订单原始数据Excel
输出：HTML报告 / Excel结果表
"""
import openpyxl, os, json, base64, shutil, hashlib, tempfile, logging
from collections import defaultdict
from datetime import datetime
from pathlib import Path

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
plt.rcParams['font.sans-serif'] = ['WenQuanYi Zen Hei','Noto Sans CJK JP','DejaVu Sans']
plt.rcParams['axes.unicode_minus'] = False

logger = logging.getLogger(__name__)

MONTHS = ['4月','5月','6月','7月','8月','9月','10月','11月','12月','1月','2月','3月']
Q_MAP = {'4月':'Q1','5月':'Q1','6月':'Q1','7月':'Q2','8月':'Q2','9月':'Q2',
         '10月':'Q3','11月':'Q3','12月':'Q3','1月':'Q4','2月':'Q4','3月':'Q4'}

class AnalysisReport:
    def __init__(self, report_id: str, org_id: str, title: str, html_path: str,
                 excel_path: str = None, ppt_path: str = None, created_at: str = None):
        self.report_id = report_id
        self.org_id = org_id
        self.title = title
        self.html_path = html_path
        self.excel_path = excel_path
        self.ppt_path = ppt_path
        self.created_at = created_at or datetime.now().strftime('%Y-%m-%d %H:%M:%S')

    def to_dict(self):
        return {
            'report_id': self.report_id,
            'org_id': self.org_id,
            'title': self.title,
            'html_path': self.html_path,
            'excel_path': self.excel_path,
            'ppt_path': self.ppt_path,
            'created_at': self.created_at,
        }


class DataAnalyzer:
    """经营数据分析引擎 - 纯本地运算，不调外部API"""

    def __init__(self, work_dir: str = None):
        self.work_dir = work_dir or tempfile.mkdtemp(prefix='analysis_')
        self.chart_dir = os.path.join(self.work_dir, 'charts')
        os.makedirs(self.chart_dir, exist_ok=True)

    def analyze(self, excel_path: str, title: str = "经营分析报告",
                org_id: str = "default") -> AnalysisReport:
        """执行全部分析，返回报告对象"""
        logger.info(f"开始分析: {excel_path}")
        start = datetime.now()

        # ─── 加载数据 ───
        wb = openpyxl.load_workbook(excel_path, read_only=True, data_only=True)
        ws = wb.active
        rows = list(ws.iter_rows(values_only=True))
        wb.close()

        data = []
        for r in rows[1:]:
            gsv_str = str(r[0]).strip() if r[0] is not None else '0'
            try: gsv = float(gsv_str) if gsv_str else 0.0
            except: gsv = 0.0
            m = str(r[34]).strip() if len(r) > 34 and r[34] else '3月'
            if m.startswith('2026'): m = '3月'
            data.append({
                'gsv': gsv,
                'brand': str(r[11] or '') if len(r) > 11 else '',
                'cat3': str(r[18] or '') if len(r) > 18 else '',  # 四级分类=实际三级分类
                'store': str(r[13] or '') if len(r) > 13 else '',
                'prov': str(r[61] or '') if len(r) > 61 else '',
                'prod_id': str(r[8] or '') if len(r) > 8 else '',
                'order_id': str(r[2] or '') if len(r) > 2 else '',
                'month': m,
            })

        logger.info(f"加载 {len(data)} 条记录, GSV={sum(d['gsv'] for d in data):.0f}")

        # ─── 聚合 ───
        TOTAL = sum(d['gsv'] for d in data)
        monthly = {m:{'gsv':0.0,'orders':set()} for m in MONTHS}
        for d in data:
            monthly[d['month']]['gsv'] += d['gsv']
            monthly[d['month']]['orders'].add(d['order_id'])

        cat3_total = defaultdict(float)
        cat3_monthly = defaultdict(lambda: defaultdict(float))
        brand_total = defaultdict(float)
        store_gsv = defaultdict(float)
        prov_total = defaultdict(float)
        prov_orders = defaultdict(set)
        prod_cat3 = defaultdict(set)
        month_prods = {m: set() for m in MONTHS}
        all_prods = set()

        for d in data:
            cat3_total[d['cat3']] += d['gsv']
            cat3_monthly[d['cat3']][d['month']] += d['gsv']
            brand_total[d['brand']] += d['gsv']
            if d['store']: store_gsv[d['store']] += d['gsv']
            if d['prov']:
                prov_total[d['prov']] += d['gsv']
                prov_orders[d['prov']].add(d['order_id'])
            if d['prod_id']:
                prod_cat3[d['cat3']].add(d['prod_id'])
                all_prods.add(d['prod_id'])
                month_prods[d['month']].add(d['prod_id'])

        cats_s = sorted(cat3_total.items(), key=lambda x: -x[1])
        brands_s = sorted(brand_total.items(), key=lambda x: -x[1])
        stores_s = sorted(store_gsv.items(), key=lambda x: -x[1])
        provs_s = sorted(prov_total.items(), key=lambda x: -x[1])

        gsv_wan = TOTAL / 10000
        ca_all = len(set(d['order_id'] for d in data))
        active_3 = len(month_prods['3月'])
        lenovo_pct = brand_total.get('联想', 0) / TOTAL
        multi_pct = brand_total.get('多品牌', 0) / TOTAL
        top3_stores = stores_s[:3]
        top3_stores_gsv = sum(v for _, v in top3_stores)
        top3_provs = provs_s[:3]
        top3_provs_gsv = sum(v for _, v in top3_provs)

        # ─── 生成图表 ───
        charts = self._generate_charts(monthly, cats_s, brand_total, brands_s,
                                       stores_s, provs_s, TOTAL)

        # ─── 生成HTML ───
        report_id = hashlib.md5(f"{org_id}{datetime.now().timestamp()}".encode()).hexdigest()[:12]
        html_path = os.path.join(self.work_dir, f'{report_id}.html')

        html = self._build_html(report_id, title, org_id, charts, data, monthly,
                                cats_s, brands_s, stores_s, provs_s,
                                TOTAL, gsv_wan, ca_all, active_3,
                                lenovo_pct, multi_pct, top3_stores, top3_stores_gsv,
                                top3_provs, top3_provs_gsv, prov_orders, prod_cat3,
                                month_prods, all_prods, brand_total, store_gsv, prov_total)

        with open(html_path, 'w', encoding='utf-8') as f:
            f.write(html)

        # ─── 生成Excel ───
        excel_path = os.path.join(self.work_dir, f'{report_id}.xlsx')
        self._generate_excel(excel_path, TOTAL, cats_s, brands_s, stores_s,
                             provs_s, brand_total, store_gsv, prov_orders, prov_total,
                             prod_cat3, month_prods, all_prods, data)

        # ─── 生成PPT ───
        ppt_path = os.path.join(self.work_dir, f'{report_id}.pptx')
        self._generate_ppt(ppt_path, title, TOTAL, cats_s, brands_s, stores_s,
                           provs_s, brand_total, store_gsv, prov_orders, prov_total,
                           prod_cat3, month_prods, all_prods, data)

        elapsed = (datetime.now() - start).total_seconds()
        logger.info(f"分析完成: {elapsed:.1f}s, report={report_id}")

        return AnalysisReport(
            report_id=report_id,
            org_id=org_id,
            title=title,
            html_path=html_path,
            excel_path=excel_path,
            ppt_path=ppt_path,
            created_at=datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
        )

    # ────────────────────────────────────────────
    # 图表生成
    # ────────────────────────────────────────────
    def _generate_charts(self, monthly, cats_s, brand_total, brands_s,
                         stores_s, provs_s, TOTAL):
        charts = {}

        # Chart 1: Monthly bar
        months_act = [m for m in MONTHS if monthly[m]['gsv'] > 0]
        vals = [monthly[m]['gsv'] / 10000 for m in months_act]
        if months_act:
            fig, ax = plt.subplots(figsize=(10, 4))
            bars = ax.bar(months_act, vals, color='#2563eb', width=0.55,
                          edgecolor='white', linewidth=0.8)
            for b, v in zip(bars, vals):
                ax.text(b.get_x() + b.get_width() / 2, b.get_height() + 0.3,
                        f'{v:.1f}', ha='center', va='bottom', fontsize=10,
                        fontweight='bold', color='#1e293b')
            ax.set_ylabel('GSV（万元）', fontsize=11, color='#64748b')
            ax.set_title('月度GSV达成', fontsize=15, fontweight='bold',
                         color='#1e293b', pad=15)
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)
            ax.spines['left'].set_color('#e2e8f0')
            ax.spines['bottom'].set_color('#e2e8f0')
            ax.tick_params(colors='#64748b', labelsize=10)
            ax.set_ylim(0, max(vals) * 1.18)
            charts['monthly'] = self._save_chart(fig, 'monthly')

        # Chart 2: Category donut
        if cats_s:
            fig, ax = plt.subplots(figsize=(7, 5))
            labels = [f'{c}  {v / TOTAL * 100:.1f}%' for c, v in cats_s]
            colors_cat = ['#2563eb', '#f59e0b', '#10b981', '#ef4444', '#8b5cf6',
                          '#ec4899', '#14b8a6', '#f97316', '#06b6d4', '#84cc16']
            wedges, texts = ax.pie([v for _, v in cats_s], labels=labels,
                                   colors=colors_cat[:len(cats_s)],
                                   startangle=90, counterclock=False,
                                   wedgeprops={'edgecolor': 'white', 'linewidth': 2.5},
                                   textprops={'fontsize': 13, 'color': '#1e293b'})
            ax.set_title('品类GSV占比', fontsize=16, fontweight='bold',
                         color='#1e293b', pad=15)
            charts['category'] = self._save_chart(fig, 'category')

        # Chart 3: Brand horizontal bar (top 10)
        brands10 = brands_s[:10]
        if brands10:
            names_b = [b for b, _ in brands10[::-1]]
            vals_b = [v / 10000 for _, v in brands10[::-1]]
            fig, ax = plt.subplots(figsize=(9, 5))
            bars = ax.barh(names_b, vals_b, color='#2563eb', height=0.5,
                           edgecolor='white', linewidth=0.5)
            for b, v in zip(bars, vals_b):
                ax.text(b.get_width() + 0.15, b.get_y() + b.get_height() / 2,
                        f'{v:.1f}万', ha='left', va='center', fontsize=10,
                        color='#64748b')
            ax.set_xlabel('GSV（万元）', fontsize=11, color='#64748b')
            ax.set_title('品牌GSV TOP10', fontsize=15, fontweight='bold',
                         color='#1e293b', pad=15)
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)
            ax.spines['left'].set_color('#e2e8f0')
            ax.spines['bottom'].set_color('#e2e8f0')
            ax.tick_params(colors='#64748b', labelsize=10)
            ax.set_xlim(0, max(vals_b) * 1.35)
            charts['brand'] = self._save_chart(fig, 'brand')

        # Chart 4: Store bar (top 10)
        stores10 = stores_s[:10]
        if stores10:
            names_s = [s[:10] + '…' if len(s) > 10 else s for s, _ in stores10[::-1]]
            vals_s = [v / 10000 for _, v in stores10[::-1]]
            fig, ax = plt.subplots(figsize=(9, 5))
            bars = ax.barh(names_s, vals_s, color='#10b981', height=0.5,
                           edgecolor='white')
            for b, v in zip(bars, vals_s):
                ax.text(b.get_width() + 0.1, b.get_y() + b.get_height() / 2,
                        f'{v:.1f}万', ha='left', va='center', fontsize=9,
                        color='#64748b')
            ax.set_xlabel('GSV（万元）', fontsize=11, color='#64748b')
            ax.set_title('商家GSV TOP10', fontsize=15, fontweight='bold',
                         color='#1e293b', pad=15)
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)
            ax.spines['left'].set_color('#e2e8f0')
            ax.spines['bottom'].set_color('#e2e8f0')
            ax.tick_params(colors='#64748b', labelsize=9)
            ax.set_xlim(0, max(vals_s) * 1.35)
            charts['store'] = self._save_chart(fig, 'store')

        # Chart 5: Province bar
        provs10 = provs_s[:10]
        if provs10:
            names_p = [p for p, _ in provs10]
            vals_p = [v / 10000 for _, v in provs10]
            colors_p = ['#2563eb', '#3b82f6', '#60a5fa', '#93c5fd', '#bfdbfe',
                        '#1d4ed8', '#1e40af', '#3730a3', '#312e81', '#1e1b4b']
            fig, ax = plt.subplots(figsize=(10, 4.5))
            bars = ax.bar(names_p, vals_p, color=colors_p[:len(provs10)],
                          width=0.55, edgecolor='white')
            for b, v in zip(bars, vals_p):
                ax.text(b.get_x() + b.get_width() / 2, b.get_height() + 0.2,
                        f'{v:.1f}', ha='center', va='bottom', fontsize=9,
                        color='#1e293b', fontweight='bold')
            ax.set_ylabel('GSV（万元）', fontsize=11, color='#64748b')
            ax.set_title('区域GSV TOP10', fontsize=15, fontweight='bold',
                         color='#1e293b', pad=15)
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)
            ax.spines['left'].set_color('#e2e8f0')
            ax.spines['bottom'].set_color('#e2e8f0')
            ax.tick_params(colors='#64748b', labelsize=9)
            plt.setp(ax.get_xticklabels(), rotation=30, ha='right', fontsize=9)
            ax.set_ylim(0, max(vals_p) * 1.2)
            charts['province'] = self._save_chart(fig, 'province')

        return charts

    def _save_chart(self, fig, name):
        path = os.path.join(self.chart_dir, f'{name}.png')
        fig.savefig(path, dpi=180, bbox_inches='tight', facecolor='white',
                    edgecolor='none')
        plt.close(fig)
        with open(path, 'rb') as f:
            b64 = base64.b64encode(f.read()).decode()
        return f'data:image/png;base64,{b64}'

    # ────────────────────────────────────────────
    # HTML生成
    # ────────────────────────────────────────────
    def _build_html(self, report_id, title, org_id, charts, data, monthly,
                    cats_s, brands_s, stores_s, provs_s,
                    TOTAL, gsv_wan, ca_all, active_3,
                    lenovo_pct, multi_pct, top3_stores, top3_stores_gsv,
                    top3_provs, top3_provs_gsv, prov_orders, prod_cat3,
                    month_prods, all_prods, brand_total, store_gsv, prov_total):

        def cat3_rows():
            rows = ''
            for i, (c, v) in enumerate(cats_s):
                bg = '#F8FAFC' if i % 2 == 0 else '#FFFFFF'
                rows += f'''<tr style="background:{bg}">
                    <td style="padding:8px 12px;font-size:13px">{c}</td>
                    <td style="padding:8px 12px;font-size:13px;text-align:right;font-weight:600">{v:,.0f}</td>
                    <td style="padding:8px 12px;font-size:13px;text-align:right">
                        <div style="display:flex;align-items:center;gap:8px;justify-content:flex-end">
                            <div style="width:80px;height:6px;background:#E2E8F0;border-radius:3px;overflow:hidden">
                                <div style="width:{v/TOTAL*100}%;height:100%;background:#2563eb;border-radius:3px"></div>
                            </div>
                            <span style="color:#64748B">{v/TOTAL*100:.1f}%</span>
                        </div>
                    </td>
                </tr>'''
            return rows

        def brand_rows():
            rows = ''
            for i, (b, v) in enumerate(brands_s[:15]):
                bg = '#F8FAFC' if i % 2 == 0 else '#FFFFFF'
                rows += f'''<tr style="background:{bg}">
                    <td style="padding:6px 10px;font-size:12px">{b}</td>
                    <td style="padding:6px 10px;font-size:12px;text-align:right">{v:,.0f}</td>
                    <td style="padding:6px 10px;font-size:12px;text-align:right">{v/TOTAL*100:.1f}%</td>
                </tr>'''
            return rows

        def store_rows():
            rows = ''
            for i, (s, v) in enumerate(stores_s[:15]):
                bg = '#F8FAFC' if i % 2 == 0 else '#FFFFFF'
                so = len([d for d in data if d['store'] == s])
                rows += f'''<tr style="background:{bg}">
                    <td style="padding:6px 10px;font-size:12px">{s}</td>
                    <td style="padding:6px 10px;font-size:12px;text-align:right">{v:,.0f}</td>
                    <td style="padding:6px 10px;font-size:12px;text-align:right">{so}</td>
                </tr>'''
            return rows

        def prov_rows():
            rows = ''
            for i, (p, v) in enumerate(provs_s):
                bg = '#F8FAFC' if i % 2 == 0 else '#FFFFFF'
                rows += f'''<tr style="background:{bg}">
                    <td style="padding:6px 10px;font-size:12px">{p}</td>
                    <td style="padding:6px 10px;font-size:12px;text-align:right">{v:,.0f}</td>
                    <td style="padding:6px 10px;font-size:12px;text-align:right">{len(prov_orders[p])}</td>
                    <td style="padding:6px 10px;font-size:12px;text-align:right">{v/TOTAL*100:.1f}%</td>
                </tr>'''
            return rows

        def prod_rows():
            rows = ''
            for c, _ in cats_s:
                pa = len(prod_cat3[c] & month_prods['3月'])
                pt = len(prod_cat3[c])
                if pa == 0: continue
                rate = pa / pt if pt > 0 else 0
                rows += f'''<tr>
                    <td style="padding:8px 12px;font-size:13px">{c}</td>
                    <td style="padding:8px 12px;font-size:13px;text-align:right;font-weight:600">{pa}</td>
                    <td style="padding:8px 12px;font-size:13px;text-align:right">{pt}</td>
                    <td style="padding:8px 12px;font-size:13px;text-align:right">
                        <div style="display:flex;align-items:center;gap:8px;justify-content:flex-end">
                            <div style="width:60px;height:6px;background:#E2E8F0;border-radius:3px;overflow:hidden">
                                <div style="width:{rate*100}%;height:100%;background:#10b981;border-radius:3px"></div>
                            </div>
                            <span style="color:#64748B">{rate:.1%}</span>
                        </div>
                    </td>
                </tr>'''
            return rows

        cat_pct1 = cats_s[0][1] / TOTAL * 100 if cats_s else 0
        cat_pct2 = cats_s[1][1] / TOTAL * 100 if len(cats_s) > 1 else 0
        lenovo_total = brand_total.get('联想', 0) + brand_total.get('联想百应', 0)
        third_party = sum(brand_total.get(b, 0) for b in ['华硕', '戴尔', '惠普'])
        top3_store_names = [s[:6] + '…' if len(s) > 6 else s for s, _ in top3_stores]
        top3_prov_names = [p for p, _ in top3_provs]

        html = f'''<!DOCTYPE html>
<html lang="zh-CN">
<head><meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>{title}</title>
<style>
  *,*::before,*::after{{box-sizing:border-box;margin:0;padding:0}}
  body{{font-family:-apple-system,BlinkMacSystemFont,"Segoe UI",Roboto,"Noto Sans SC","PingFang SC","Microsoft YaHei",sans-serif;background:#F1F5F9;color:#1E293B;line-height:1.6}}
  .header{{background:linear-gradient(135deg,#0F172A 0%,#1E293B 100%);color:#fff;padding:40px 0 36px;border-bottom:4px solid #2563EB}}
  .header .inner{{max-width:1100px;margin:0 auto;padding:0 32px}}
  .header h1{{font-size:26px;font-weight:700;margin-bottom:4px}}
  .header .sub{{color:#94A3B8;font-size:14px}}
  .header .meta{{color:#64748B;font-size:12px;margin-top:10px;display:flex;gap:20px}}
  .kpi-bar{{max-width:1100px;margin:-24px auto 0;padding:0 32px;display:grid;grid-template-columns:repeat(4,1fr);gap:14px;position:relative;z-index:10}}
  .kpi-card{{background:#fff;border-radius:10px;padding:18px 22px;box-shadow:0 1px 3px rgba(0,0,0,.06);border-top:3px solid #2563EB}}
  .kpi-card:nth-child(2){{border-top-color:#10B981}}
  .kpi-card:nth-child(3){{border-top-color:#F59E0B}}
  .kpi-card:nth-child(4){{border-top-color:#8B5CF6}}
  .kpi-card .label{{font-size:11px;color:#64748B;margin-bottom:3px}}
  .kpi-card .value{{font-size:24px;font-weight:700;color:#1E293B}}
  .kpi-card .unit{{font-size:12px;color:#94A3B8;font-weight:400;margin-left:2px}}
  .container{{max-width:1100px;margin:0 auto;padding:0 32px}}
  .section{{margin-top:28px}}
  .section-title{{font-size:17px;font-weight:700;color:#1E293B;padding-bottom:8px;border-bottom:2px solid #E2E8F0;margin-bottom:16px;display:flex;align-items:center;gap:10px}}
  .section-title .badge{{font-size:10px;font-weight:500;padding:2px 10px;border-radius:20px;background:#DBEAFE;color:#2563EB}}
  .card-grid{{display:grid;gap:18px}}
  .card-grid.two{{grid-template-columns:1fr 1fr}}
  .card{{background:#fff;border-radius:10px;box-shadow:0 1px 3px rgba(0,0,0,.05);overflow:hidden}}
  .card .body{{padding:18px}}
  .card .ctitle{{font-size:13px;font-weight:600;color:#64748B;margin-bottom:12px;letter-spacing:.02em}}
  .card img{{width:100%;height:auto;display:block}}
  .data-table{{width:100%;border-collapse:collapse}}
  .data-table th{{background:#1E293B;color:#fff;font-size:11px;font-weight:600;padding:8px 10px;text-align:left;white-space:nowrap}}
  .data-table th.r{{text-align:right}}
  .data-table tr:hover{{background:#F1F5F9!important}}
  .full-card{{background:#fff;border-radius:10px;box-shadow:0 1px 3px rgba(0,0,0,.05);padding:22px}}
  .insight{{background:#EFF6FF;border-left:4px solid #2563EB;border-radius:0 8px 8px 0;padding:12px 16px;margin-top:14px;font-size:13px;color:#1E40AF;line-height:1.7}}
  .insight.green{{background:#F0FDF4;border-left-color:#10B981;color:#166534}}
  .insight.purple{{background:#FAF5FF;border-left-color:#8B5CF6;color:#5B21B6}}
  .ccl-grid{{display:grid;grid-template-columns:1fr 1fr;gap:14px}}
  .ccl-item{{background:#fff;border-radius:10px;padding:16px 18px;box-shadow:0 1px 3px rgba(0,0,0,.05);border-left:3px solid #2563EB}}
  .ccl-item .emoji{{font-size:18px;margin-right:6px}}
  .ccl-item .t{{font-weight:600;font-size:13px;margin-bottom:3px}}
  .ccl-item .d{{font-size:11px;color:#64748B}}
  .ccl-item:nth-child(2){{border-left-color:#10B981}}
  .ccl-item:nth-child(3){{border-left-color:#F59E0B}}
  .ccl-item:nth-child(4){{border-left-color:#8B5CF6}}
  .ccl-item:nth-child(5){{border-left-color:#EF4444}}
  .ccl-item:nth-child(6){{border-left-color:#EC4899}}
  .footer{{text-align:center;padding:28px;color:#94A3B8;font-size:11px;margin-top:36px;border-top:1px solid #E2E8F0}}
  .safe-badge{{display:inline-flex;align-items:center;gap:6px;background:#ECFDF5;color:#065F46;font-size:11px;padding:4px 12px;border-radius:20px;font-weight:500;margin-left:10px}}
  @media(max-width:768px){{.kpi-bar{{grid-template-columns:1fr 1fr}}.card-grid.two{{grid-template-columns:1fr}}.ccl-grid{{grid-template-columns:1fr}}}}
</style></head>
<body>
<div class="header">
  <div class="inner">
    <h1>📊 经营分析报告 <span class="safe-badge">🔒 数据安全</span></h1>
    <div class="sub">{title}</div>
    <div class="meta">
      <span>📅 {datetime.now().strftime("%Y年%m月%d日")}</span>
      <span>🏢 {org_id}</span>
      <span>🔬 纯本地运算 · 数据不存留</span>
    </div>
  </div>
</div>

<div class="kpi-bar">
  <div class="kpi-card"><div class="label">GSV</div><div class="value">{gsv_wan:.0f}<span class="unit">万元</span></div></div>
  <div class="kpi-card"><div class="label">订单量</div><div class="value">{ca_all}<span class="unit">单</span></div></div>
  <div class="kpi-card"><div class="label">动销商品</div><div class="value">{active_3}<span class="unit">个</span></div></div>
  <div class="kpi-card"><div class="label">活跃品牌</div><div class="value">{len(brand_total)}<span class="unit">个</span></div></div>
</div>

<div class="container">
  <div class="section">
    <div class="section-title">📈 整体达成 <span class="badge">月度快照</span></div>
    <div class="card"><div class="body"><img src="{charts.get('monthly','')}" alt="月度GSV"></div></div>
    <div class="insight"><strong>3月GSV达成{gsv_wan:.0f}万</strong>，订单量{ca_all}单，动销商品{active_3}个。</div>
  </div>

  <div class="section">
    <div class="section-title">📦 品类维度 <span class="badge">{len(cats_s)}个品类</span></div>
    <div class="card-grid two">
      <div class="card"><div class="body"><img src="{charts.get('category','')}" alt="品类占比"></div></div>
      <div class="card"><div class="body"><div class="ctitle">品类GSV明细</div>
      <table class="data-table"><thead><tr><th>品类</th><th class="r">GSV</th><th class="r">占比</th></tr></thead><tbody>{cat3_rows()}</tbody></table></div></div>
    </div>
    <div class="insight"><strong>核心发现：</strong>Top1 {cats_s[0][0]}占比{cat_pct1:.1f}%；Top2 {cats_s[1][0] if len(cats_s)>1 else ''}占比{cat_pct2:.1f}%。</div>
  </div>

  <div class="section">
    <div class="section-title">🏷️ 品牌维度 <span class="badge">{len(brand_total)}个品牌</span></div>
    <div class="card-grid two">
      <div class="card"><div class="body"><img src="{charts.get('brand','')}" alt="品牌TOP10"></div></div>
      <div class="card"><div class="body"><div class="ctitle">品牌GSV TOP15</div>
      <table class="data-table"><thead><tr><th>品牌</th><th class="r">GSV</th><th class="r">占比</th></tr></thead><tbody>{brand_rows()}</tbody></table></div></div>
    </div>
    <div class="insight green"><strong>品牌洞察：</strong>联想系（含联想百应）占比{lenovo_total/TOTAL*100:.1f}%，多品牌{multi_pct*100:.1f}%；第三方品牌（华硕/戴尔/惠普）合计占比{third_party/TOTAL*100:.1f}%。</div>
  </div>

  <div class="section">
    <div class="section-title">🛒 商品维度 <span class="badge">动销率 {active_3/len(all_prods)*100:.1f}%</span></div>
    <div class="full-card">
      <div style="font-size:13px;color:#64748B;margin-bottom:10px">3月动销 <strong>{active_3}</strong>个 / 上架 <strong>{len(all_prods)}</strong>个</div>
      <table class="data-table"><thead><tr><th>品类</th><th class="r">动销数</th><th class="r">上架数</th><th class="r">动销率</th></tr></thead><tbody>{prod_rows()}</tbody></table>
    </div>
    <div class="insight"><strong>商品表现：</strong>3月整体动销率{active_3/len(all_prods)*100:.1f}%（{active_3}/{len(all_prods)}）。</div>
  </div>

  <div class="section">
    <div class="section-title">🏪 商家维度 <span class="badge">{len(store_gsv)}家店铺</span></div>
    <div class="card-grid two">
      <div class="card"><div class="body"><img src="{charts.get('store','')}" alt="商家TOP10"></div></div>
      <div class="card"><div class="body"><div class="ctitle">商家GSV TOP15</div>
      <table class="data-table"><thead><tr><th>店铺</th><th class="r">GSV</th><th class="r">订单数</th></tr></thead><tbody>{store_rows()}</tbody></table></div></div>
    </div>
    <div class="insight purple"><strong>商家格局：</strong>Top3（{'/'.join(top3_store_names)}）合计GSV {top3_stores_gsv:,.0f}，占比{top3_stores_gsv/TOTAL*100:.1f}%。</div>
  </div>

  <div class="section">
    <div class="section-title">📍 区域维度 <span class="badge">{len(prov_total)}个省份</span></div>
    <div class="card-grid two">
      <div class="card"><div class="body"><img src="{charts.get('province','')}" alt="区域TOP10"></div></div>
      <div class="card"><div class="body"><div class="ctitle">区域GSV明细</div>
      <table class="data-table"><thead><tr><th>省份</th><th class="r">GSV</th><th class="r">订单量</th><th class="r">占比</th></tr></thead><tbody>{prov_rows()}</tbody></table></div></div>
    </div>
    <div class="insight green"><strong>区域分布：</strong>Top3（{'/'.join(top3_prov_names)}）合计占比{top3_provs_gsv/TOTAL*100:.1f}%。</div>
  </div>

  <div class="section">
    <div class="section-title">🎯 核心结论</div>
    <div class="ccl-grid">
      <div class="ccl-item"><div class="t"><span class="emoji">📊</span>3月业绩</div><div class="d">GSV {gsv_wan:.0f}万，订单{ca_all}单，动销{active_3}个，品牌{len(brand_total)}个</div></div>
      <div class="ccl-item"><div class="t"><span class="emoji">📦</span>品类聚焦</div><div class="d">{cats_s[0][0]} GSV {cats_s[0][1]:,.0f}，占比{cat_pct1:.1f}%</div></div>
      <div class="ccl-item"><div class="t"><span class="emoji">🏷️</span>品牌策略</div><div class="d">联想系占比{lenovo_total/TOTAL*100:.1f}%，多品牌{multi_pct*100:.1f}%</div></div>
      <div class="ccl-item"><div class="t"><span class="emoji">🏪</span>商家生态</div><div class="d">{len(store_gsv)}家店铺，Top3占比{top3_stores_gsv/TOTAL*100:.1f}%</div></div>
      <div class="ccl-item"><div class="t"><span class="emoji">📍</span>区域分布</div><div class="d">{len(prov_total)}个省份，Top3合计{top3_provs_gsv/TOTAL*100:.1f}%</div></div>
      <div class="ccl-item"><div class="t"><span class="emoji">🔒</span>数据安全</div><div class="d">原始数据已自动删除，报告仅您可见</div></div>
    </div>
  </div>
</div>
<div class="footer">百应智星 · 经营分析数字员工 · 数据本地化运算</div>
</body></html>'''
        return html

    def _generate_excel(self, excel_path, TOTAL, cats_s, brands_s, stores_s,
                        provs_s, brand_total, store_gsv, prov_orders, prov_total,
                        prod_cat3, month_prods, all_prods, data):
        """生成Excel分析结果"""
        wb = openpyxl.Workbook()

        # Sheet 1: 整体达成
        ws1 = wb.active
        ws1.title = '整体达成'
        ws1.append(['指标', '数值'])
        ws1.append(['GSV（万元）', round(TOTAL / 10000, 2)])
        ws1.append(['订单量（单）', len(set(d['order_id'] for d in data))])
        ws1.append(['动销商品（个）', len(month_prods.get('3月', set()))])
        ws1.append(['上架商品（个）', len(all_prods)])
        ws1.append(['活跃品牌（个）', len(brand_total)])

        # Sheet 2: 品类维度
        ws2 = wb.create_sheet('品类维度')
        ws2.append(['品类', 'GSV', '占比(%)'])
        for c, v in cats_s:
            ws2.append([c, round(v, 2), round(v/TOTAL*100, 1)])

        # Sheet 3: 品牌维度
        ws3 = wb.create_sheet('品牌维度')
        ws3.append(['品牌', 'GSV', '占比(%)'])
        for b, v in brands_s:
            ws3.append([b, round(v, 2), round(v/TOTAL*100, 1)])

        # Sheet 4: 商家维度
        ws4 = wb.create_sheet('商家维度')
        ws4.append(['店铺', 'GSV', '订单数'])
        for s, v in stores_s:
            order_n = len([d for d in data if d['store'] == s])
            ws4.append([s, round(v, 2), order_n])

        # Sheet 5: 区域维度
        ws5 = wb.create_sheet('区域维度')
        ws5.append(['省份', 'GSV', '订单量', '占比(%)'])
        for p, v in provs_s:
            ws5.append([p, round(v, 2), len(prov_orders[p]), round(v/TOTAL*100, 1)])

        # Sheet 6: 商品维度
        ws6 = wb.create_sheet('商品维度')
        ws6.append(['品类', '动销数', '上架数', '动销率(%)'])
        for c, _ in cats_s:
            pa = len(prod_cat3[c] & month_prods['3月'])
            pt = len(prod_cat3[c])
            if pa > 0:
                ws6.append([c, pa, pt, round(pa/pt*100, 1) if pt else 0])

        out = excel_path
        wb.save(out)
        logger.info(f"Excel报告已生成: {out}")
        return out

    def _generate_ppt(self, ppt_path, title, TOTAL, cats_s, brands_s, stores_s,
                      provs_s, brand_total, store_gsv, prov_orders, prov_total,
                      prod_cat3, month_prods, all_prods, data):
        """生成PPT分析报告"""
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        def add_slide(title_text, content_fn):
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
            # Title
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = title_text
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0x1E, 0x29, 0x3B)
            # Content
            content_fn(slide)
            return slide

        # Slide 1: Cover
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0x0F, 0x17, 0x2A)
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.2))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER
        txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11), Inches(0.6))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = f"生成时间: {datetime.now().strftime('%Y-%m-%d')}"
        p2.font.size = Pt(16)
        p2.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
        p2.alignment = PP_ALIGN.CENTER

        # Slide 2: KPI
        def kpi_slide(slide):
            kpis = [
                (f"GSV: {TOTAL/10000:.0f}万元", f"订单: {len(set(d['order_id'] for d in data))}单"),
                (f"动销: {len(month_prods.get('3月', set()))}个", f"品牌: {len(brand_total)}个"),
            ]
            for i, (k1, k2) in enumerate(kpis):
                txBox = slide.shapes.add_textbox(Inches(0.5 + i*6.5), Inches(1.2), Inches(6), Inches(1.5))
                tf = txBox.text_frame
                p = tf.paragraphs[0]
                p.text = k1
                p.font.size = Pt(20)
                p.font.color.rgb = RGBColor(0x1E, 0x29, 0x3B)
                p2 = tf.add_paragraph()
                p2.text = k2
                p2.font.size = Pt(16)
                p2.font.color.rgb = RGBColor(0x64, 0x74, 0x8B)
        add_slide('📊 整体达成', kpi_slide)

        # Slide 3: Category
        def cat_slide(slide):
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(0.5))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = f"共{len(cats_s)}个品类  |  Top1: {cats_s[0][0]} {cats_s[0][1]/TOTAL*100:.1f}%"
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(0x64, 0x74, 0x8B)
            table = slide.shapes.add_table(len(cats_s)+1, 3, Inches(0.5), Inches(2), Inches(12), Inches(0.4*(len(cats_s)+1))).table
            for j, h in enumerate(['品类', 'GSV', '占比(%)']):
                cell = table.cell(0, j)
                cell.text = h
                for p in cell.text_frame.paragraphs:
                    p.font.bold = True
                    p.font.size = Pt(12)
            for i, (c, v) in enumerate(cats_s):
                table.cell(i+1, 0).text = c
                table.cell(i+1, 1).text = f"{v:,.0f}"
                table.cell(i+1, 2).text = f"{v/TOTAL*100:.1f}%"
        add_slide('📦 品类维度', cat_slide)

        # Slide 4: Brand
        def brand_slide(slide):
            lenovo_total = brand_total.get('联想', 0) + brand_total.get('联想百应', 0)
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(0.5))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = f"共{len(brand_total)}个品牌  |  联想系占比{lenovo_total/TOTAL*100:.1f}%"
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(0x64, 0x74, 0x8B)
            rows_n = min(20, len(brands_s))
            table = slide.shapes.add_table(rows_n+1, 3, Inches(0.5), Inches(2), Inches(12), Inches(0.35*(rows_n+1))).table
            for j, h in enumerate(['品牌', 'GSV', '占比(%)']):
                cell = table.cell(0, j)
                cell.text = h
                for p in cell.text_frame.paragraphs:
                    p.font.bold = True
                    p.font.size = Pt(11)
            for i, (b, v) in enumerate(brands_s[:rows_n]):
                table.cell(i+1, 0).text = b
                table.cell(i+1, 1).text = f"{v:,.0f}"
                table.cell(i+1, 2).text = f"{v/TOTAL*100:.1f}%"
        add_slide('🏷️ 品牌维度', brand_slide)

        # Slide 5: Store
        def store_slide(slide):
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(0.5))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            top3 = sum(v for _, v in stores_s[:3])
            p.text = f"共{len(store_gsv)}家店铺  |  Top3占比{top3/TOTAL*100:.1f}%"
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(0x64, 0x74, 0x8B)
            rows_n = min(20, len(stores_s))
            table = slide.shapes.add_table(rows_n+1, 3, Inches(0.5), Inches(2), Inches(12), Inches(0.35*(rows_n+1))).table
            for j, h in enumerate(['店铺', 'GSV', '占比(%)']):
                cell = table.cell(0, j)
                cell.text = h
                for p in cell.text_frame.paragraphs:
                    p.font.bold = True
                    p.font.size = Pt(11)
            for i, (s, v) in enumerate(stores_s[:rows_n]):
                table.cell(i+1, 0).text = s
                table.cell(i+1, 1).text = f"{v:,.0f}"
                table.cell(i+1, 2).text = f"{v/TOTAL*100:.1f}%"
        add_slide('🏪 商家维度', store_slide)

        # Slide 6: Province
        def prov_slide(slide):
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(0.5))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            top3 = sum(v for _, v in provs_s[:3])
            p.text = f"共{len(prov_total)}个省份  |  Top3占比{top3/TOTAL*100:.1f}%"
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(0x64, 0x74, 0x8B)
            table = slide.shapes.add_table(len(provs_s)+1, 4, Inches(0.5), Inches(2), Inches(12), Inches(0.35*(len(provs_s)+1))).table
            for j, h in enumerate(['省份', 'GSV', '订单量', '占比(%)']):
                cell = table.cell(0, j)
                cell.text = h
                for p in cell.text_frame.paragraphs:
                    p.font.bold = True
                    p.font.size = Pt(11)
            for i, (p, v) in enumerate(provs_s):
                table.cell(i+1, 0).text = p
                table.cell(i+1, 1).text = f"{v:,.0f}"
                table.cell(i+1, 2).text = f"{len(prov_orders[p])}"
                table.cell(i+1, 3).text = f"{v/TOTAL*100:.1f}%"
        add_slide('📍 区域维度', prov_slide)

        # Slide 7: Product
        def prod_slide(slide):
            active_3 = len(month_prods.get('3月', set()))
            all_n = len(all_prods)
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(0.5))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = f"动销{active_3}个 / 上架{all_n}个  |  整体动销率{active_3/all_n*100:.1f}%"
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(0x64, 0x74, 0x8B)
            prod_items = [(c, len(prod_cat3[c] & month_prods['3月']), len(prod_cat3[c])) for c, _ in cats_s if len(prod_cat3[c] & month_prods['3月']) > 0]
            table = slide.shapes.add_table(len(prod_items)+1, 4, Inches(0.5), Inches(2), Inches(12), Inches(0.35*(len(prod_items)+1))).table
            for j, h in enumerate(['品类', '动销数', '上架数', '动销率(%)']):
                cell = table.cell(0, j)
                cell.text = h
                for p in cell.text_frame.paragraphs:
                    p.font.bold = True
                    p.font.size = Pt(11)
            for i, (c, pa, pt) in enumerate(prod_items):
                table.cell(i+1, 0).text = c
                table.cell(i+1, 1).text = str(pa)
                table.cell(i+1, 2).text = str(pt)
                table.cell(i+1, 3).text = f"{pa/pt*100:.1f}%" if pt else "0%"
        add_slide('🛒 商品维度', prod_slide)

        prs.save(ppt_path)
        logger.info(f"PPT报告已生成: {ppt_path}")
        return ppt_path

    def cleanup(self):
        """删除工作目录（原始数据和分析产物）"""
        if os.path.exists(self.work_dir):
            shutil.rmtree(self.work_dir)
            logger.info(f"已清理: {self.work_dir}")
