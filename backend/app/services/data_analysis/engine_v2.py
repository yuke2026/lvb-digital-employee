"""
经营分析 V2 — 多表灵活分析引擎
支持：1-5个表上传 → 关联合并 → 多维度组合分析 → 多章节报告(HTML/PPT)
"""
import os, json, base64, hashlib, logging, tempfile
from datetime import datetime
from pathlib import Path
from io import BytesIO
from collections import defaultdict

import pandas as pd
import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
plt.rcParams['font.sans-serif'] = ['WenQuanYi Zen Hei','Noto Sans CJK JP','DejaVu Sans']
plt.rcParams['axes.unicode_minus'] = False

logger = logging.getLogger(__name__)

# ─── 色板 ───
CLR = ['#2563eb','#10b981','#f59e0b','#ef4444','#8b5cf6',
       '#ec4899','#14b8a6','#f97316','#06b6d4','#84cc16',
       '#6366f1','#d946ef','#0ea5e9','#22c55e','#eab308',
       '#a855f7','#3b82f6','#34d399','#fb923c','#f472b6']

def _color_map(values):
    seen = {}
    for i, v in enumerate(values):
        if v not in seen:
            seen[v] = CLR[i % len(CLR)]
    return seen


class MultiTableAnalyzer:
    """多表灵活分析引擎"""

    def __init__(self, work_dir: str = None):
        self.work_dir = work_dir or tempfile.mkdtemp(prefix='analysis_v2_')
        self.chart_dir = os.path.join(self.work_dir, 'charts')
        os.makedirs(self.chart_dir, exist_ok=True)
        self.tables = {}        # table_id → DataFrame
        self.meta = {}          # table_id → {columns, dtypes, rows}

    # ────────────────────────────────────────────
    # 加载
    # ────────────────────────────────────────────
    def load_table(self, file_path: str, table_id: str = None) -> dict:
        """加载一个Excel/CSV文件，返回schema信息"""
        ext = Path(file_path).suffix.lower()
        if ext == '.csv':
            df = pd.read_csv(file_path)
        else:
            df = pd.read_excel(file_path, engine='openpyxl')

        # 清洗列名
        df.columns = [str(c).strip() for c in df.columns]
        # 去掉全空列
        df = df.dropna(axis=1, how='all')

        tid = table_id or hashlib.md5(file_path.encode()).hexdigest()[:8]
        self.tables[tid] = df

        schema = self._schema(df)
        self.meta[tid] = schema
        logger.info(f"加载表 {tid}: {schema['cols']}列 × {len(df)}行")
        return {**schema, 'table_id': tid, 'rows': len(df)}

    def _schema(self, df: pd.DataFrame) -> dict:
        """提取schema信息"""
        cols = []
        for c in df.columns:
            dtype = str(df[c].dtype)
            if pd.api.types.is_numeric_dtype(df[c]):
                vtype = '数值'
            elif pd.api.types.is_datetime64_any_dtype(df[c]):
                vtype = '日期'
            else:
                try:
                    pd.to_datetime(df[c], format='mixed')
                    vtype = '日期'
                except:
                    vtype = '文本'
            cols.append({
                'name': c,
                'dtype': dtype,
                'type': vtype,
                'non_null': int(df[c].notna().sum()),
                'unique': int(df[c].nunique()),
            })
        return {'cols': len(cols), 'columns': cols}

    def preview(self, table_id: str, n: int = 20) -> list:
        """返回前n行数据"""
        if table_id not in self.tables:
            raise ValueError(f"表 {table_id} 不存在")
        df = self.tables[table_id].head(n)
        return json.loads(df.to_json(orient='records', force_ascii=False))

    # ────────────────────────────────────────────
    # 关联合并
    # ────────────────────────────────────────────
    def merge(self, joins: list) -> pd.DataFrame:
        """按关联规则合并多个表"""
        if len(self.tables) < 1:
            raise ValueError("至少需要1个表")

        ids = list(self.tables.keys())
        result = self.tables[ids[0]].copy()

        for j in joins:
            lt = j['left_table']
            lc = j['left_col']
            rt = j['right_table']
            rc = j['right_col']
            how = j.get('how', 'left')

            if lt not in self.tables or rt not in self.tables:
                raise ValueError(f"表不存在: {lt}/{rt}")

            left_df = self.tables[lt]
            right_df = self.tables[rt]
            if lc in left_df.columns and rc in right_df.columns:
                if left_df[lc].dtype != right_df[rc].dtype:
                    try:
                        right_df[rc] = right_df[rc].astype(str)
                        left_df[lc] = left_df[lc].astype(str)
                    except:
                        pass

            result = result.merge(
                right_df, how=how,
                left_on=lc, right_on=rc,
                suffixes=('', f'_{rt}')
            )

        logger.info(f"合并完成: {len(result)}行 × {len(result.columns)}列")
        return result

    # ────────────────────────────────────────────
    # 多维度组合分析
    # ────────────────────────────────────────────
    def analyze_all(self, df: pd.DataFrame,
                    row_dims: list = None,
                    values: list = None,
                    filters: list = None,
                    dimensions: list = None) -> dict:
        """多维度组合分析
        
        两种模式：
        1. 命名维度模式 (dimensions) — 逐个添加命名分析维度，每维度独立配置
           dimensions: [
               {
                   name: '商家维度',
                   row_dims: [{column, label}],
                   values: [{column, agg, label}],
                   filters: [{column, op, value}],
               },
               ...
           ]
        2. 交叉模式 (row_dims × values) — 所有行维×所有值维自动组合
        
        row_dims: [{column, label}] — 行维度列表（交叉模式）
        values: [{column, agg, label}] — 值维度列表（交叉模式）
        filters: [{column, op, value}] — 全局过滤器
        """
        row_dims = row_dims or []
        values = values or []
        filters = filters or []
        filters_copy = [dict(f) for f in filters]  # 避免修改原对象

        sections = []

        if dimensions:
            # ── 命名维度模式 ──
            for dim in dimensions:
                name = dim.get('name', '分析')
                d_row_dims = dim.get('row_dims', [])
                d_values = dim.get('values', [])
                d_filters = dim.get('filters', []) or filters_copy

                # 应用该维度的过滤器
                filtered = df.copy()
                for f in d_filters:
                    col, op, val = f['column'], f.get('op', 'eq'), f.get('value')
                    if col not in filtered.columns:
                        continue
                    if op == 'eq':       filtered = filtered[filtered[col] == val]
                    elif op == 'neq':    filtered = filtered[filtered[col] != val]
                    elif op == 'gt':     filtered = filtered[filtered[col] > float(val)]
                    elif op == 'gte':    filtered = filtered[filtered[col] >= float(val)]
                    elif op == 'lt':     filtered = filtered[filtered[col] < float(val)]
                    elif op == 'lte':    filtered = filtered[filtered[col] <= float(val)]
                    elif op == 'in':
                        vals = val.split(',') if isinstance(val, str) else val
                        filtered = filtered[filtered[col].isin(vals)]

                # 生成该维度的分析章节
                if d_row_dims and d_values:
                    for rd in d_row_dims:
                        for v in d_values:
                            sec = self._single_section(filtered, [rd['column']], rd['label'], v)
                            sec['dim_name'] = name
                            sections.append(sec)
                elif d_row_dims and not d_values:
                    # 只有行维度，用计数
                    sec = self._single_section(filtered, [d_row_dims[0]['column']], d_row_dims[0]['label'],
                                               {'column': '', 'agg': 'count', 'label': '计数'})
                    sec['dim_name'] = name
                    sections.append(sec)
                elif d_values and not d_row_dims:
                    # 只有值维度，无分组汇总
                    sec = self._single_section(filtered, [], '总计', d_values[0])
                    sec['dim_name'] = name
                    if not sec.get('charts'):
                        sec['charts'] = {}
                    sections.append(sec)
                else:
                    # 纯计数（整体分析）
                    total = len(filtered)
                    sections.append({
                        'type': 'overall', 'dim_name': name,
                        'label': name,
                        'row_dim': {'column': '', 'label': '总计'},
                        'value': {'column': '', 'agg': 'count', 'label': '总行数'},
                        'columns': ['总计', '行数'],
                        'data': [{'总计': name, '行数': total}],
                        'rows': 1, 'charts': {}, 'total': total,
                        'agg_label': '计数',
                    })

            dims_meta = {'named': [d.get('name','') for d in dimensions]}

        else:
            # ── 交叉模式（向后兼容） ──
            for f in filters_copy:
                col, op, val = f['column'], f.get('op', 'eq'), f.get('value')
                if col not in df.columns:
                    continue
                if op == 'eq':       df = df[df[col] == val]
                elif op == 'neq':    df = df[df[col] != val]
                elif op == 'gt':     df = df[df[col] > float(val)]
                elif op == 'gte':    df = df[df[col] >= float(val)]
                elif op == 'lt':     df = df[df[col] < float(val)]
                elif op == 'lte':    df = df[df[col] <= float(val)]
                elif op == 'in':
                    vals = val.split(',') if isinstance(val, str) else val
                    df = df[df[col].isin(vals)]

            for rd in row_dims:
                for v in values:
                    sec = self._single_section(df, [rd['column']], rd['label'], v)
                    sec['dim_name'] = rd['label']
                    sections.append(sec)

            if len(row_dims) > 1 and values:
                for v in values:
                    combined_cols = [d['column'] for d in row_dims]
                    combined_display = '+'.join(d['label'][:6] for d in row_dims)
                    sec = self._single_section(df, combined_cols, combined_display, v)
                    sec['type'] = 'combined'
                    sec['dim_name'] = combined_display
                    sections.append(sec)

            if not sections:
                total = len(df)
                sections.append({
                    'type': 'overall', 'dim_name': '总览',
                    'label': '总览',
                    'row_dim': {'column': '', 'label': '计数'},
                    'value': {'column': '', 'agg': 'count', 'label': '总行数'},
                    'columns': ['计数'],
                    'data': [{'计数': total}],
                    'rows': 1, 'charts': {}, 'total': total,
                    'agg_label': '计数',
                })

            dims_meta = {'row': row_dims, 'values': values}

        return {
            'sections': sections,
            'dimensions': dims_meta,
            'total_rows': len(df),
        }

    def _single_section(self, df, group_cols, row_label, value_conf):
        """计算单个 (行维度组合 × 值维度) 的分析"""
        col = value_conf['column']
        agg = value_conf.get('agg', 'sum')
        val_label = value_conf.get('label', col)
        agg_labels = {'sum': '求和', 'avg': '平均', 'count': '计数',
                      'min': '最小值', 'max': '最大值'}

        # 确保数值类型
        if col and col in df.columns and agg in ('sum', 'avg', 'mean', 'min', 'max'):
            df[col] = pd.to_numeric(df[col], errors='coerce').fillna(0)

        if group_cols and col and col in df.columns:
            agg_result = df.groupby(group_cols).agg({col: agg}).reset_index()
        elif group_cols:
            agg_result = df.groupby(group_cols).size().reset_index(name='计数')
        elif col and col in df.columns:
            agg_result = df.agg({col: agg}).to_frame().T
        else:
            agg_result = pd.DataFrame({'计数': [len(df)]})

        # 排序
        val_col = col if col in agg_result.columns else '计数'
        agg_result = agg_result.sort_values(val_col, ascending=False)

        # 生成图表
        chart_col = col if col else '计数'
        charts = self._single_chart(agg_result, group_cols, chart_col, val_label, row_label)

        # 转数据
        columns = [str(c) for c in agg_result.columns]
        data = json.loads(agg_result.to_json(orient='records', force_ascii=False))
        total_val = sum(r.get(chart_col, 0) for r in data) if data else 0

        return {
            'type': 'single',
            'label': f"{row_label} × {val_label}",
            'row_dim': {'column': group_cols[0] if group_cols else '', 'label': row_label},
            'value': value_conf,
            'columns': columns,
            'data': data,
            'rows': len(data),
            'charts': charts,
            'total': total_val,
            'agg_label': agg_labels.get(agg, agg),
        }

    def _single_chart(self, df, group_cols, val_col, val_label, row_label):
        """为单个分析生成图表"""
        charts = {}
        if len(df) == 0:
            return charts

        if not group_cols:
            return charts

        # 获取分组标签列（如果多个分组列，拼接它们）
        if len(group_cols) == 1:
            name_col = group_cols[0]
            names = [str(v)[:15] for v in df[name_col].values]
        else:
            # 多列组合 → 拼接标签
            names = []
            for _, row in df[group_cols].iterrows():
                label = ' / '.join(str(row[c])[:10] for c in group_cols)
                names.append(label)

        vals = pd.to_numeric(df[val_col], errors='coerce').fillna(0).values

        if len(names) == 0:
            return charts

        # 柱状图
        df_plot = df.copy()
        df_plot['_label'] = names
        df_plot['_val'] = vals
        df_plot = df_plot.sort_values('_val', ascending=False).head(20)

        names_plot = df_plot['_label'].tolist()
        vals_plot = df_plot['_val'].tolist()

        fig, ax = plt.subplots(figsize=(max(8, len(names_plot)*0.35), 4.5))
        cmap = _color_map(names_plot)
        colors = [cmap[n] for n in names_plot]

        if len(names_plot) <= 8:
            bars = ax.bar(names_plot, vals_plot, color=colors, width=0.55,
                          edgecolor='white', linewidth=0.8)
            for b, v in zip(bars, vals_plot):
                ax.text(b.get_x() + b.get_width()/2, b.get_height(),
                        f'{v:,.0f}', ha='center', va='bottom',
                        fontsize=9, fontweight='bold', color='#1e293b')
            plt.setp(ax.get_xticklabels(), rotation=25, ha='right', fontsize=8)
        else:
            ax.barh(names_plot[::-1], vals_plot[::-1], color=colors[::-1],
                    height=0.5, edgecolor='white')
            for i, (b, v) in enumerate(zip(ax.patches, vals_plot[::-1])):
                ax.text(b.get_width() + max(vals_plot)*0.01,
                        b.get_y() + b.get_height()/2,
                        f'{v:,.0f}', ha='left', va='center', fontsize=8,
                        color='#64748b')

        ax.set_title(f'{row_label} × {val_label}', fontsize=14, fontweight='bold',
                     color='#1e293b', pad=12)
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        ax.spines['left'].set_color('#e2e8f0')
        ax.spines['bottom'].set_color('#e2e8f0')
        ax.tick_params(colors='#64748b', labelsize=9)
        charts['main'] = self._save_chart(fig, f'main_{hashlib.md5(row_label.encode()).hexdigest()[:8]}')

        # 饼图（最多12个）
        if len(names_plot) <= 12:
            fig2, ax2 = plt.subplots(figsize=(7, 5))
            total = sum(vals_plot)
            labels_pie = [f'{n}  {v/total*100:.1f}%' for n, v in zip(names_plot[:10], vals_plot[:10])]
            if len(vals_plot) > 10:
                other = sum(vals_plot[10:])
                labels_pie.append(f'其他  {other/total*100:.1f}%')
                vals_pie = list(vals_plot[:10]) + [other]
            else:
                vals_pie = list(vals_plot[:10])
            cmap_pie = _color_map(labels_pie)
            colors_pie = [cmap_pie[l] for l in labels_pie]
            ax2.pie(vals_pie, labels=labels_pie, colors=colors_pie,
                    startangle=90, counterclock=False,
                    wedgeprops={'edgecolor': 'white', 'linewidth': 2},
                    textprops={'fontsize': 11, 'color': '#1e293b'})
            ax2.set_title(f'{row_label} × {val_label} 占比', fontsize=15,
                          fontweight='bold', color='#1e293b', pad=15)
            charts['pie'] = self._save_chart(fig2, f'pie_{hashlib.md5(row_label.encode()).hexdigest()[:8]}')

        return charts

    # ────────────────────────────────────────────
    # 报告生成
    # ────────────────────────────────────────────
    def build_html_report(self, title: str, analysis_result: dict,
                          join_info: list = None, org_id: str = 'default',
                          ai_insights: str = '') -> str:
        """生成多章节HTML报告"""
        sections = analysis_result.get('sections', [])
        dims = analysis_result.get('dimensions', {})
        total_rows = analysis_result.get('total_rows', 0)
        now_str = datetime.now().strftime("%Y年%m月%d日 %H:%M")

        # 构建每个分析章节的HTML
        sections_html = ''
        for i, sec in enumerate(sections):
            sec_id = f"sec_{i}"
            label = sec.get('label', '分析')
            total = sec.get('total', 0)
            agg_label = sec.get('agg_label', '')
            row_count = sec.get('rows', 0)
            charts = sec.get('charts', {})
            data = sec.get('data', [])
            cols = sec.get('columns', [])

            # 图表
            charts_html = ''
            if 'main' in charts:
                charts_html += f'<div class="card"><div class="body"><img src="{charts["main"]}" alt="{label}"></div></div>'
            if 'pie' in charts:
                charts_html += f'<div style="margin-top:18px" class="card"><div class="body"><img src="{charts["pie"]}" alt="{label}占比"></div></div>'

            # 表格行
            table_rows = ''
            for j, r in enumerate(data[:50]):
                bg = '#F8FAFC' if j % 2 == 0 else '#FFFFFF'
                cells = ''.join(
                    f'<td style="padding:5px 8px;font-size:11px;text-align:{"right" if c == cols[-1] else "left"}">'
                    f'{r.get(c, "")}</td>' for c in cols[:6]
                )
                table_rows += f'<tr style="background:{bg}">{cells}</tr>'

            # 统计卡
            kpi_cards = ''
            if total > 0:
                kpi_cards += f'''
                <div class="kpi-mini"><span class="kpi-mini-label">{agg_label}</span>
                <span class="kpi-mini-val">{total:,.0f}</span></div>'''
            kpi_cards += f'''
                <div class="kpi-mini"><span class="kpi-mini-label">分组数</span>
                <span class="kpi-mini-val">{row_count}</span></div>'''

            sec_type_badge = '组合分析' if sec.get('type') == 'combined' else '维度分析'

            sections_html += f'''
            <div class="section" id="{sec_id}">
              <div class="section-title">{i+1}. {label} <span class="badge">{sec_type_badge}</span></div>
              <div class="kpi-row-mini">{kpi_cards}</div>
              {charts_html}
              <div style="margin-top:16px">
                <div class="section-subtitle">📋 明细数据 <span class="badge">{min(len(data), 50)}/{len(data)}行</span></div>
                <div class="full-card">'''
            if table_rows:
                sections_html += '<table class="data-table"><thead><tr>' + \
                    ''.join(f'<th>{" " if c == cols[-1] else ""}{c}</th>' for c in cols[:6]) + \
                    '</tr></thead><tbody>' + table_rows + '</tbody></table>'
                if len(data) > 50:
                    sections_html += f'<p style="text-align:center;padding:6px;font-size:11px;color:#94a3b8">仅显示前50行（共{len(data)}行）</p>'
            else:
                sections_html += '<p style="color:#94a3b8;font-size:13px;text-align:center;padding:20px">暂无数据</p>'
            sections_html += '</div></div></div>'

        # AI洞察
        insights_html = ''
        if ai_insights:
            insights_html = f'''
            <div class="section">
              <div class="section-title">🧠 AI 洞察 <span class="badge">智能分析</span></div>
              <div class="insight-card">{ai_insights}</div>
            </div>'''

        # 维度标签
        row_labels = ''.join(
            f'<span class="dim-tag row">{d["label"]}</span>'
            for d in dims.get('row', [])
        )
        val_labels = ''.join(
            f'<span class="dim-tag val">{v.get("label", v["column"])}</span>'
            for v in dims.get('values', [])
        )

        report_id = hashlib.md5(f"{org_id}{datetime.now().timestamp()}".encode()).hexdigest()[:12]

        html = f'''<!DOCTYPE html>
<html lang="zh-CN">
<head><meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>{title}</title>
<style>
  *,*::before,*::after{{box-sizing:border-box;margin:0;padding:0}}
  body{{font-family:-apple-system,BlinkMacSystemFont,"Segoe UI",Roboto,"Noto Sans SC","PingFang SC","Microsoft YaHei",sans-serif;background:#F1F5F9;color:#1E293B;line-height:1.6}}
  .header{{background:linear-gradient(135deg,#0F172A 0%,#1E293B 100%);color:#fff;padding:40px 0 36px;border-bottom:4px solid #2563EB}}
  .header .inner{{max-width:1100px;margin:0 auto;padding:0 32px}}
  .header h1{{font-size:26px;font-weight:700;margin-bottom:4px}}
  .header .sub{{color:#94A3B8;font-size:14px}}
  .header .meta{{color:#64748B;font-size:12px;margin-top:10px;display:flex;gap:20px}}
  .kpi-row{{max-width:1100px;margin:-24px auto 0;padding:0 32px;display:grid;grid-template-columns:repeat(auto-fit, minmax(150px,1fr));gap:12px;position:relative;z-index:10}}
  .kpi-card{{background:#fff;border-radius:10px;padding:14px 18px;box-shadow:0 1px 3px rgba(0,0,0,.06);border-top:3px solid #2563EB}}
  .kpi-card .label{{font-size:11px;color:#64748B;margin-bottom:2px}}
  .kpi-card .value{{font-size:18px;font-weight:700;color:#1E293B}}
  .kpi-row-mini{{display:flex;gap:8px;margin-bottom:14px;flex-wrap:wrap}}
  .kpi-mini{{background:#F8FAFC;border:1px solid #E2E8F0;border-radius:8px;padding:8px 14px;display:flex;gap:8px;align-items:center}}
  .kpi-mini-label{{font-size:11px;color:#64748B}}
  .kpi-mini-val{{font-size:16px;font-weight:700;color:#1E293B}}
  .container{{max-width:1100px;margin:0 auto;padding:0 32px}}
  .section{{margin-top:28px;padding-top:8px;border-top:2px solid #E2E8F0}}
  .section:first-of-type{{border-top:none;margin-top:36px}}
  .section-title{{font-size:17px;font-weight:700;color:#1E293B;padding-bottom:8px;border-bottom:2px solid #E2E8F0;margin-bottom:16px;display:flex;align-items:center;gap:10px}}
  .section-subtitle{{font-size:14px;font-weight:600;color:#475569;padding-bottom:6px;margin-bottom:10px;display:flex;align-items:center;gap:8px}}
  .badge{{font-size:10px;font-weight:500;padding:2px 10px;border-radius:20px;background:#DBEAFE;color:#2563EB;white-space:nowrap}}
  .card{{background:#fff;border-radius:10px;box-shadow:0 1px 3px rgba(0,0,0,.05);overflow:hidden;margin-bottom:14px}}
  .card .body{{padding:16px}}
  .card img{{width:100%;height:auto;display:block}}
  .data-table{{width:100%;border-collapse:collapse}}
  .data-table th{{background:#1E293B;color:#fff;font-size:11px;font-weight:600;padding:6px 8px;text-align:left;white-space:nowrap}}
  .data-table tr:hover{{background:#F1F5F9!important}}
  .full-card{{background:#fff;border-radius:10px;box-shadow:0 1px 3px rgba(0,0,0,.05);padding:16px;overflow-x:auto}}
  .dim-tag{{display:inline-block;padding:2px 10px;border-radius:12px;font-size:11px;font-weight:500;margin:2px}}
  .dim-tag.row{{background:#DBEAFE;color:#2563EB}}
  .dim-tag.val{{background:#DCFCE7;color:#16A34A}}
  .insight-card{{background:#FFF7ED;border:1px solid #FED7AA;border-radius:10px;padding:20px;font-size:14px;line-height:1.8;color:#7C2D12}}
  .insight-card h3,.insight-card strong{{color:#C2410C}}
  .insight-card ul{{margin:8px 0;padding-left:18px}}
  .insight-card li{{margin:4px 0}}
  .insight-card hr{{border:none;border-top:1px solid #FED7AA;margin:12px 0}}
  .footer{{text-align:center;padding:28px;color:#94A3B8;font-size:11px;margin-top:36px;border-top:1px solid #E2E8F0}}
  @media(max-width:768px){{.kpi-row{{grid-template-columns:1fr 1fr}}}}
</style></head>
<body>
<div class="header">
  <div class="inner">
    <h1>📊 {title}</h1>
    <div class="sub">多维度综合分析报告</div>
    <div class="meta">
      <span>📅 {now_str}</span>
      <span>🔬 纯本地运算 · 数据安全</span>
      <span>📊 {len(sections)} 个分析维度</span>
    </div>
  </div>
</div>

<div class="kpi-row">
  <div class="kpi-card"><div class="label">分析维度</div><div class="value">{len(dims.get("row", []))}<span style="font-size:12px;color:#94a3b8;font-weight:400"> 个行维度</span></div></div>
  <div class="kpi-card"><div class="label">分析指标</div><div class="value">{len(dims.get("values", []))}<span style="font-size:12px;color:#94a3b8;font-weight:400"> 个值维度</span></div></div>
  <div class="kpi-card"><div class="label">分析章节</div><div class="value">{len(sections)}<span style="font-size:12px;color:#94a3b8;font-weight:400"> 个</span></div></div>
  <div class="kpi-card"><div class="label">数据总量</div><div class="value">{total_rows:,}<span style="font-size:12px;color:#94a3b8;font-weight:400"> 行</span></div></div>
</div>

<div class="container">
  {row_labels}{val_labels}

  {sections_html}

  {insights_html}
</div>

<div class="footer">
  百应智星 · 经营分析 V2 · 纯本地运算 · 由AI驱动
</div>
</body>
</html>'''
        return html, report_id

    # ────────────────────────────────────────────
    # PPT生成
    # ────────────────────────────────────────────
    def build_ppt_report(self, title: str, analysis_result: dict,
                         ai_insights: str = '', org_id: str = 'default') -> bytes:
        """生成PPT报告（python-pptx）"""
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        sections = analysis_result.get('sections', [])
        dims = analysis_result.get('dimensions', {})
        now_str = datetime.now().strftime("%Y-%m-%d %H:%M")

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # ── 配色 ──
        DARK = RGBColor(0x0F, 0x17, 0x2A)
        BLUE = RGBColor(0x25, 0x63, 0xEB)
        GRAY = RGBColor(0x94, 0x94, 0xB8)
        WHITE = RGBColor(0xFF, 0xFF, 0xFF)

        def _add_bg(slide, color=DARK):
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = color

        def _add_textbox(slide, left, top, width, height, text,
                         font_size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
            txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                              Inches(width), Inches(height))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(font_size)
            p.font.bold = bold
            p.font.color.rgb = color
            p.alignment = align
            return tf

        def _add_table_slide(prs, sec_idx, sec, section_title):
            """为每个分析章节生成一页PPT"""
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
            _add_bg(slide, DARK)

            # 标题
            _add_textbox(slide, 0.5, 0.3, 12, 0.6,
                         f"{sec_idx+1}. {section_title}",
                         font_size=24, bold=True, color=WHITE)

            data = sec.get('data', [])
            cols = sec.get('columns', [])[:6]
            if not data or not cols:
                _add_textbox(slide, 0.5, 1.5, 12, 0.5,
                             "暂无数据", font_size=14, color=GRAY)
                return

            # 表格
            rows = min(len(data), 20) + 1  # +1 表头
            cols_count = len(cols)
            table = slide.shapes.add_table(rows, cols_count,
                                           Inches(0.5), Inches(1.2),
                                           Inches(12.3), Inches(0.4 * rows)).table

            # 表头
            for ci, c in enumerate(cols):
                cell = table.cell(0, ci)
                cell.text = c
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(11)
                    p.font.bold = True
                    p.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = BLUE

            # 数据行
            for ri, row_data in enumerate(data[:rows-1]):
                for ci, c in enumerate(cols):
                    cell = table.cell(ri + 1, ci)
                    cell.text = str(row_data.get(c, ''))
                    for p in cell.text_frame.paragraphs:
                        p.font.size = Pt(10)
                        p.font.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B) if ri % 2 == 0 else RGBColor(0x0F, 0x17, 0x2A)

            # Summary
            total = sec.get('total', 0)
            agg_label = sec.get('agg_label', '')
            _add_textbox(slide, 0.5, 6.5, 12, 0.5,
                         f"{agg_label}: {total:,.0f}  |  分组: {sec.get('rows', 0)}  |  数据: {len(data)}行",
                         font_size=12, color=GRAY)

        # ── 封面页 ──
        slide1 = prs.slides.add_slide(prs.slide_layouts[6])
        _add_bg(slide1, DARK)
        _add_textbox(slide1, 0.5, 1.5, 12, 1.5, "📊 经营分析报告",
                     font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _add_textbox(slide1, 0.5, 3.2, 12, 1, title,
                     font_size=24, color=GRAY, align=PP_ALIGN.CENTER)
        _add_textbox(slide1, 0.5, 4.5, 12, 0.5,
                     f"生成时间: {now_str}",
                     font_size=14, color=GRAY, align=PP_ALIGN.CENTER)
        _add_textbox(slide1, 0.5, 5.0, 12, 0.5,
                     f"分析维度: {len(sections)} 个分析章节",
                     font_size=14, color=GRAY, align=PP_ALIGN.CENTER)
        _add_textbox(slide1, 0.5, 6.5, 12, 0.5,
                     "百应智星 · 纯本地运算",
                     font_size=12, color=GRAY, align=PP_ALIGN.CENTER)

        # ── 概述页 ──
        slide2 = prs.slides.add_slide(prs.slide_layouts[6])
        _add_bg(slide2, DARK)
        _add_textbox(slide2, 0.5, 0.3, 12, 0.6, "📊 分析概览",
                     font_size=28, bold=True, color=WHITE)
        row_labels = ', '.join(d['label'] for d in dims.get('row', []))
        val_labels = ', '.join(v.get('label', v['column']) for v in dims.get('values', []))
        _add_textbox(slide2, 0.5, 1.5, 12, 0.5,
                     f"行维度: {row_labels}", font_size=16, color=GRAY)
        _add_textbox(slide2, 0.5, 2.2, 12, 0.5,
                     f"值维度: {val_labels}", font_size=16, color=GRAY)
        _add_textbox(slide2, 0.5, 2.9, 12, 0.5,
                     f"分析章节: {len(sections)}", font_size=16, color=GRAY)

        # 章节列表
        for i, sec in enumerate(sections):
            _add_textbox(slide2, 0.5, 3.8 + i * 0.5, 12, 0.5,
                         f"  {i+1}. {sec.get('label', '-')}",
                         font_size=14, color=WHITE)

        # ── AI洞察页 ──
        if ai_insights:
            slide_ai = prs.slides.add_slide(prs.slide_layouts[6])
            _add_bg(slide_ai, DARK)
            _add_textbox(slide_ai, 0.5, 0.3, 12, 0.6, "🧠 AI 洞察分析",
                         font_size=28, bold=True, color=WHITE)
            # 分段显示
            clean_text = ai_insights.replace('<h3>', '\n').replace('</h3>', '\n')
            clean_text = clean_text.replace('<strong>', '').replace('</strong>', '')
            clean_text = clean_text.replace('<ul>', '').replace('</ul>', '')
            clean_text = clean_text.replace('<li>', '• ').replace('</li>', '\n')
            clean_text = clean_text.replace('<hr>', '\n')
            clean_text = clean_text.replace('<p>', '').replace('</p>', '\n')
            clean_text = clean_text.replace('<br>', '\n')
            _add_textbox(slide_ai, 0.5, 1.2, 12, 5.5, clean_text,
                         font_size=13, color=GRAY)

        # ── 每个分析章节 ──
        for i, sec in enumerate(sections):
            _add_table_slide(prs, i, sec, sec.get('label', f'分析{i+1}'))

        # 尾页
        slide_end = prs.slides.add_slide(prs.slide_layouts[6])
        _add_bg(slide_end, DARK)
        _add_textbox(slide_end, 0.5, 2.5, 12, 1,
                     "谢谢", font_size=44, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER)
        _add_textbox(slide_end, 0.5, 4, 12, 0.5,
                     "百应智星 · 经营分析 V2",
                     font_size=18, color=GRAY, align=PP_ALIGN.CENTER)

        buf = BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf.getvalue()

    def _save_chart(self, fig, name):
        path = os.path.join(self.chart_dir, f'{name}.png')
        fig.savefig(path, dpi=180, bbox_inches='tight', facecolor='white',
                    edgecolor='none')
        plt.close(fig)
        with open(path, 'rb') as f:
            b64 = base64.b64encode(f.read()).decode()
        return f'data:image/png;base64,{b64}'

    def cleanup(self):
        """清理临时文件"""
        import shutil
        if os.path.exists(self.work_dir):
            shutil.rmtree(self.work_dir, ignore_errors=True)
